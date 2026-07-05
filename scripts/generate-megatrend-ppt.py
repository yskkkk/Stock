#!/usr/bin/env python3
"""Professional mega-trend investment deck with charts, TOC, and section dividers."""

from __future__ import annotations

import tempfile
from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import font_manager
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "reports" / "megatrend-investment-report-2026.pptx"
OUT_FALLBACK = OUT.with_name("megatrend-investment-report-2026-v3.pptx")
CHART_DIR = Path(tempfile.mkdtemp(prefix="megatrend-charts-"))

# ── Palette ──────────────────────────────────────────────────────────
NAVY = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
DARK = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_LINE = RGBColor(0xE2, 0xE8, 0xF0)

MPL_COLORS = ["#06B6D4", "#8B5CF6", "#10B981", "#F59E0B", "#F43F5E", "#3B82F6", "#EC4899", "#14B8A6"]
TREND_COLORS = {
    "AI 인프라": "#06B6D4",
    "원자력": "#10B981",
    "반도체": "#8B5CF6",
    "양자": "#F59E0B",
    "로보틱스": "#F43F5E",
    "우주": "#3B82F6",
    "바이오": "#EC4899",
    "경제안보": "#14B8A6",
}

_slide_no = 0

# 티커 → 한글 회사명 (발표 자료 가독성)
TICKER_KO: dict[str, str] = {
    "NVDA": "엔비디아", "AVGO": "브로드컴", "MRVL": "마벨", "AMD": "AMD",
    "INTC": "인텔", "TSM": "TSMC", "MU": "마이크론", "GFS": "글로벌파운드리",
    "AMAT": "어플라이드 머티리얼즈", "LRCX": "램리서치", "KLAC": "KLA", "ASML": "ASML",
    "ARM": "ARM 홀딩스", "ON": "온세미", "MPWR": "모노리틱 파워", "AMKR": "암코",
    "VRT": "버티브", "ETN": "이튼", "PWR": "콴타 서비스", "CAT": "캐터필러",
    "GNRC": "제네락", "STX": "씨게이트", "ABB": "ABB", "EQIX": "에퀴닉스",
    "DLR": "디지털 리얼티", "ANET": "아리스타", "SMCI": "슈퍼마이크로",
    "CEG": "컨스텔레이션 에너지", "VST": "비스트라", "TLN": "탈렌 에너지",
    "D": "도미니언 에너지", "NEE": "넥스트에라", "OKLO": "오클로", "SMR": "뉴스케일",
    "NNE": "나노 뉴클리어", "CCJ": "카메코", "BWXT": "BWX 테크놀로지",
    "LEU": "센트러스 에너지", "GEV": "GE 버노바", "BE": "블룸에너지",
    "IONQ": "아이온큐", "RGTI": "리게티", "QBTS": "디웨이브", "QUBT": "퀀텀컴퓨팅",
    "IBM": "IBM", "GOOGL": "알파벳(구글)", "MSFT": "마이크로소프트", "AMZN": "아마존",
    "META": "메타", "PANW": "팔로알토", "CRWD": "크라우드스트라이크",
    "NET": "클라우드플레어", "ZS": "지스케일러", "LAES": "실스큐", "WKEY": "와이즈키",
    "NXPI": "NXP", "TSLA": "테슬라", "ISRG": "인튜이티브 서지컬",
    "TER": "테라다인", "CGNX": "코그넥스", "ROK": "로크웰", "FANUY": "파낙",
    "SYM": "심보틱", "AVAV": "에어로바이론먼트", "KTOS": "크라토스",
    "SPCX": "스페이스X", "RKLB": "로켓랩", "FLY": "파이어플라이",
    "ASTS": "AST 스페이스모바일", "IRDM": "이리듐", "GSAT": "글로벌스타",
    "PL": "플래닛랩스", "BKSY": "블랙스카이", "LUNR": "인튜이티브 머신스",
    "LMT": "록히드마틴", "NOC": "노스롭 그루먼", "RTX": "RTX", "LHX": "L3해리스",
    "GD": "제너럴다이내믹스", "HII": "헌팅턴 잉걸스",
    "LLY": "일라이 릴리", "NVO": "노보 노디스크", "AMGN": "암젠", "VKTX": "바이킹",
    "GPCR": "스트럭처", "TMO": "써모피셔", "DHR": "다나허", "ILMN": "일루미나",
    "CRSP": "크리스퍼", "DXCM": "덱스컴", "MP": "MP 머티리얼즈",
    "USAR": "USA 레어어스", "FCX": "프리포트 맥모란", "SCCO": "서던 코퍼",
    "PLTR": "팔란티어", "BHP": "BHP", "RIO": "리오틴토", "TECK": "텍",
    "ALB": "알버말", "HUT": "헛 8",
    "SOXX": "반도체 ETF", "NUKZ": "원자력 ETF", "ARKX": "우주 ETF",
    "BOTZ": "로봇 ETF", "HACK": "사이버 ETF", "XBI": "바이오 ETF",
}
NAME_KO: dict[str, str] = {
    "Intel": "INTC · 인텔", "TSMC": "TSM · TSMC", "Samsung": "삼성전자",
    "Micron": "MU · 마이크론", "Google": "GOOGL · 알파벳", "Alphabet": "GOOGL · 알파벳",
    "Microsoft": "MSFT · 마이크로소프트", "Amazon": "AMZN · 아마존",
    "Vertiv": "VRT · 버티브", "Eaton": "ETN · 이튼", "Caterpillar": "CAT · 캐터필러",
    "IonQ": "IONQ · 아이온큐", "Rocket Lab": "RKLB · 로켓랩", "SpaceX": "SPCX · 스페이스X",
    "Figure AI": "피겨 AI(비상장)", "Tesla": "TSLA · 테슬라", "Nvidia": "NVDA · 엔비디아",
}


def t(sym: str) -> str:
    s = sym.strip().upper()
    ko = TICKER_KO.get(s)
    return f"{s} · {ko}" if ko else sym.strip()


def ts(*syms: str, sep: str = "   ") -> str:
    return sep.join(t(s) for s in syms)


def ts_block(*syms: str) -> str:
    return "\n".join(t(s) for s in syms)


def ts_cell(raw: str) -> str:
    import re
    parts = re.split(r"[\s·,/]+", raw.strip())
    out = []
    for p in parts:
        if not p:
            continue
        u = p.upper()
        if u in TICKER_KO:
            out.append(t(u))
        elif p in NAME_KO:
            out.append(NAME_KO[p])
        else:
            out.append(p)
    return "\n".join(out)


def annotate(text: str) -> str:
    import re
    for tick in sorted(TICKER_KO, key=len, reverse=True):
        ko = TICKER_KO[tick]
        text = re.sub(rf"\b{tick}\b", f"{tick}({ko})", text)
    for name, labeled in NAME_KO.items():
        text = text.replace(name, labeled)
    return text


def _korean_font() -> str:
    for name in ("Malgun Gothic", "AppleGothic", "NanumGothic", "DejaVu Sans"):
        if name in {f.name for f in font_manager.fontManager.ttflist}:
            return name
    return "DejaVu Sans"


FONT = _korean_font()
plt.rcParams.update({
    "font.family": FONT,
    "axes.unicode_minus": False,
    "figure.facecolor": "#F8FAFC",
    "axes.facecolor": "#FFFFFF",
    "axes.edgecolor": "#E2E8F0",
    "axes.labelcolor": "#334155",
    "xtick.color": "#64748B",
    "ytick.color": "#64748B",
    "grid.color": "#E2E8F0",
    "text.color": "#0F172A",
})


def _save_chart(name: str) -> Path:
    path = CHART_DIR / f"{name}.png"
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight", facecolor=plt.gcf().get_facecolor())
    plt.close()
    return path


# ── Chart builders (matplotlib) ───────────────────────────────────────

def chart_hyperscaler_capex() -> Path:
    years = ["2024", "2025", "2026", "2027"]
    values = [250, 400, 600, 1000]
    fig, ax = plt.subplots(figsize=(9, 4.5))
    bars = ax.bar(years, values, color=MPL_COLORS[0], width=0.55, edgecolor="white", linewidth=1.5)
    for bar, v in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 18,
                f"${v}B", ha="center", va="bottom", fontsize=12, fontweight="bold", color="#0F172A")
    ax.set_ylabel("Combined Capex (USD Billion)", fontsize=11)
    ax.set_title("하이퍼스케일러 AI Capex 추이 (MS · Goldman · Company Guidance)", fontsize=14, fontweight="bold", pad=14)
    ax.set_ylim(0, 1150)
    ax.yaxis.grid(True, linestyle="--", alpha=0.7)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart("capex")


def chart_capex_allocation() -> Path:
    labels = ["물리 인프라\n(전력·냉각·건설)", "반도체·GPU", "네트워킹·기타"]
    sizes = [78, 17, 5]
    colors = ["#06B6D4", "#8B5CF6", "#94A3B8"]
    fig, ax = plt.subplots(figsize=(8, 5))
    wedges, texts, autotexts = ax.pie(
        sizes, labels=labels, autopct="%1.0f%%", startangle=90,
        colors=colors, explode=(0.04, 0, 0),
        textprops={"fontsize": 11},
        wedgeprops={"edgecolor": "white", "linewidth": 2},
    )
    for t in autotexts:
        t.set_fontweight("bold")
        t.set_fontsize(13)
    ax.set_title("AI 데이터센터 Capex 배분\n(칩 ≠ 전체의 대부분)", fontsize=14, fontweight="bold", pad=16)
    return _save_chart("allocation")


def chart_dc_power() -> Path:
    years = ["2024", "2026", "2030", "2035"]
    pct = [4.4, 6.5, 9.0, 12.0]
    fig, ax = plt.subplots(figsize=(9, 4.5))
    ax.fill_between(years, pct, alpha=0.25, color=MPL_COLORS[2])
    ax.plot(years, pct, marker="o", linewidth=3, markersize=10, color=MPL_COLORS[2])
    for x, y in zip(years, pct):
        ax.annotate(f"{y}%", (x, y), textcoords="offset points", xytext=(0, 12),
                    ha="center", fontsize=11, fontweight="bold")
    ax.set_ylabel("미국 전력 대비 DC 비중 (%)", fontsize=11)
    ax.set_title("데이터센터 전력 수요 전망 (Morningstar · EEI)", fontsize=14, fontweight="bold", pad=14)
    ax.set_ylim(0, 14)
    ax.yaxis.grid(True, linestyle="--", alpha=0.7)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart("dc_power")


def chart_megatrend_market_2030() -> Path:
    trends = ["AI Infra", "GLP-1", "Utilities", "Quantum", "Space", "Humanoid", "CHIPS"]
    values = [1000, 175, 280, 11.5, 80, 12, 53]
    colors = [MPL_COLORS[i % len(MPL_COLORS)] for i in range(len(trends))]
    fig, ax = plt.subplots(figsize=(10, 5))
    bars = ax.barh(trends, values, color=colors, height=0.6, edgecolor="white")
    ax.set_xlabel("시장 규모 추정 ($B, 로그 스케일 근사)", fontsize=11)
    ax.set_title("2030 전후 메가트렌드 시장 규모 비교", fontsize=14, fontweight="bold", pad=14)
    for bar, v in zip(bars, values):
        label = f"${v}B" if v >= 10 else f"${v}B"
        ax.text(bar.get_width() + 8, bar.get_y() + bar.get_height() / 2,
                label, va="center", fontsize=10, fontweight="bold")
    ax.set_xlim(0, 1150)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart("market_2030")


def chart_demand_heatmap() -> Path:
    trends = list(TREND_COLORS.keys())
    periods = ["2025-26", "2027-28", "2029-30", "2031-35"]
    # intensity 0-100
    data = [
        [95, 90, 75, 55],  # AI
        [80, 85, 70, 90],  # nuclear
        [90, 85, 70, 50],  # semi
        [40, 65, 85, 90],  # quantum
        [35, 55, 75, 85],  # robotics
        [50, 65, 80, 90],  # space
        [95, 85, 70, 55],  # bio
        [85, 80, 75, 70],  # security
    ]
    fig, ax = plt.subplots(figsize=(10, 5.5))
    im = ax.imshow(data, cmap="YlGnBu", aspect="auto", vmin=0, vmax=100)
    ax.set_xticks(range(len(periods)))
    ax.set_xticklabels(periods, fontsize=10)
    ax.set_yticks(range(len(trends)))
    ax.set_yticklabels(trends, fontsize=10)
    for i in range(len(trends)):
        for j in range(len(periods)):
            ax.text(j, i, data[i][j], ha="center", va="center", fontsize=9,
                    color="white" if data[i][j] > 55 else "#334155", fontweight="bold")
    ax.set_title("메가트렌드별 수요 강도 Heatmap (구조적 추정)", fontsize=14, fontweight="bold", pad=14)
    cbar = plt.colorbar(im, ax=ax, fraction=0.03, pad=0.02)
    cbar.set_label("수요 강도", fontsize=10)
    return _save_chart("heatmap")


def chart_glp1_share() -> Path:
    labels = ["LLY · 일라이 릴리", "NVO · 노보 노디스크", "기타"]
    sizes = [62, 31, 7]
    colors = ["#06B6D4", "#8B5CF6", "#CBD5E1"]
    fig, ax = plt.subplots(figsize=(7, 5))
    wedges, texts, autotexts = ax.pie(sizes, labels=labels, autopct="%1.0f%%", colors=colors,
                                       startangle=140, wedgeprops={"edgecolor": "white", "linewidth": 2},
                                       textprops={"fontsize": 11})
    for t in autotexts:
        t.set_fontweight("bold")
    ax.set_title("GLP-1 시장 점유율 전망 (2030, TD Cowen)", fontsize=14, fontweight="bold", pad=14)
    return _save_chart("glp1")


def chart_quantum_growth() -> Path:
    years = ["2025", "2026", "2028", "2030", "2032"]
    market = [2.0, 2.9, 6.1, 11.5, 19.0]
    fig, ax = plt.subplots(figsize=(9, 4.5))
    ax.plot(years, market, marker="s", linewidth=3, markersize=9, color=MPL_COLORS[3])
    ax.fill_between(range(len(years)), market, alpha=0.15, color=MPL_COLORS[3])
    ax.set_xticks(range(len(years)))
    ax.set_xticklabels(years)
    for i, v in enumerate(market):
        ax.annotate(f"${v}B", (i, v), textcoords="offset points", xytext=(0, 10), ha="center", fontweight="bold")
    ax.set_ylabel("글로벌 시장 ($B)", fontsize=11)
    ax.set_title("양자컴퓨팅 시장 성장 전망 (Alora Advisory)", fontsize=14, fontweight="bold", pad=14)
    ax.yaxis.grid(True, linestyle="--", alpha=0.7)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart("quantum")


def chart_ionq_vs_google() -> Path:
    categories = ["게이트\n충실도", "연결성", "QEC\n실증", "상용\n매출", "확장\n잠재력"]
    ionq = [99, 95, 55, 90, 70]
    google = [88, 65, 98, 40, 92]
    x = range(len(categories))
    w = 0.35
    fig, ax = plt.subplots(figsize=(9, 5))
    ax.bar([i - w / 2 for i in x], ionq, w, label="IonQ · 아이온큐", color="#06B6D4", edgecolor="white")
    ax.bar([i + w / 2 for i in x], google, w, label="Google · 알파벳", color="#8B5CF6", edgecolor="white")
    ax.set_xticks(list(x))
    ax.set_xticklabels(categories, fontsize=10)
    ax.set_ylim(0, 110)
    ax.set_ylabel("상대 점수 (0–100)", fontsize=11)
    ax.set_title("IonQ vs Google Quantum AI — 기술 포지셔닝", fontsize=14, fontweight="bold", pad=14)
    ax.legend(loc="upper right", framealpha=0.9)
    ax.yaxis.grid(True, linestyle="--", alpha=0.7)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart("ionq_google")


def chart_nuclear_ppa() -> Path:
    companies = ["Meta\n패키지", "Constellation\n(CEG)", "Vistra\n(VST)", "Oklo\n(SMR)", "TerraPower\n(SMR)"]
    gw = [6.6, 1.1, 2.1, 1.2, 2.0]
    types = ["Total", "Operating", "Operating", "Future", "Future"]
    type_colors = {"Total": "#0F172A", "Operating": "#10B981", "Future": "#F59E0B"}
    colors = [type_colors[t] for t in types]
    fig, ax = plt.subplots(figsize=(9, 4.5))
    bars = ax.bar(companies, gw, color=colors, width=0.55, edgecolor="white")
    for bar, v in zip(bars, gw):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.08,
                f"{v} GW", ha="center", fontweight="bold", fontsize=11)
    ax.set_ylabel("전력 용량 (GW)", fontsize=11)
    ax.set_title("빅테크 원자력 PPA · Meta 6.6GW 패키지", fontsize=14, fontweight="bold", pad=14)
    legend = [mpatches.Patch(color="#10B981", label="가동 원전"), mpatches.Patch(color="#F59E0B", label="SMR (2030–35)")]
    ax.legend(handles=legend, loc="upper right")
    ax.set_ylim(0, 8)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart("nuclear")


def chart_value_chain() -> Path:
    fig, ax = plt.subplots(figsize=(10, 4))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 4)
    ax.axis("off")
    boxes = [
        (0.3, 2.2, "Hyperscaler\nCapex $600B+", "#0F172A", "white"),
        (2.5, 3.0, "1차: 칩\n엔비디아·브로드컴·마벨", "#8B5CF6", "white"),
        (2.5, 1.2, "2차: 물리 인프라\n버티브·이튼·콴타·캐터필러", "#06B6D4", "white"),
        (5.5, 2.2, "3차: 전력\n컨스텔레이션·비스트라", "#10B981", "white"),
        (8.0, 2.2, "AI 서비스\nMSFT·META·GOOGL", "#F59E0B", "white"),
    ]
    for x, y, text, bg, fg in boxes:
        rect = mpatches.FancyBboxPatch((x, y), 1.8, 1.0, boxstyle="round,pad=0.08",
                                        facecolor=bg, edgecolor="white", linewidth=2)
        ax.add_patch(rect)
        ax.text(x + 0.9, y + 0.5, text, ha="center", va="center", fontsize=9,
                color=fg, fontweight="bold")
    arrows = [(2.1, 2.7, 2.5, 3.4), (2.1, 2.7, 2.5, 1.7), (4.3, 3.4, 5.5, 2.9),
              (4.3, 1.7, 5.5, 2.5), (7.3, 2.7, 8.0, 2.7)]
    for x1, y1, x2, y2 in arrows:
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="->", color="#94A3B8", lw=2))
    ax.text(5, 0.3, "← 2026 수익 집중: 2차 물리 인프라 레이어", ha="center", fontsize=12,
            fontweight="bold", color="#06B6D4")
    ax.set_title("AI Value Chain — 자금 흐름", fontsize=14, fontweight="bold", pad=8)
    return _save_chart("value_chain")


def chart_horizon() -> Path:
    horizons = ["0–2년\n(Now)", "3–5년", "5–10년"]
    ai = [95, 60, 30]
    nuclear = [75, 85, 70]
    quantum = [35, 80, 95]
    bio = [90, 60, 40]
    x = range(len(horizons))
    w = 0.2
    fig, ax = plt.subplots(figsize=(9, 5))
    ax.bar([i - 1.5 * w for i in x], ai, w, label="AI 인프라", color=MPL_COLORS[0])
    ax.bar([i - 0.5 * w for i in x], nuclear, w, label="원자력", color=MPL_COLORS[2])
    ax.bar([i + 0.5 * w for i in x], quantum, w, label="양자", color=MPL_COLORS[3])
    ax.bar([i + 1.5 * w for i in x], bio, w, label="GLP-1", color=MPL_COLORS[6])
    ax.set_xticks(list(x))
    ax.set_xticklabels(horizons, fontsize=11)
    ax.set_ylabel("투자 매력도 (상대)", fontsize=11)
    ax.set_title("투자 Horizon별 섹터 매력도", fontsize=14, fontweight="bold", pad=14)
    ax.legend(loc="upper right", ncol=2, fontsize=9)
    ax.set_ylim(0, 110)
    ax.yaxis.grid(True, linestyle="--", alpha=0.7)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart("horizon")


def chart_stock_count() -> Path:
    trends = list(TREND_COLORS.keys())
    counts = [22, 21, 24, 22, 22, 22, 24, 24]
    colors = list(TREND_COLORS.values())
    fig, ax = plt.subplots(figsize=(10, 5))
    bars = ax.bar(trends, counts, color=colors, edgecolor="white", linewidth=1.5)
    ax.set_ylabel("분석 종목 수", fontsize=11)
    ax.set_title("트렌드별 커버리지 종목 수 (각 20+)", fontsize=14, fontweight="bold", pad=14)
    plt.xticks(rotation=25, ha="right", fontsize=9)
    for bar in bars:
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3,
                f"{int(bar.get_height())}", ha="center", fontweight="bold")
    ax.set_ylim(0, 28)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart("stock_count")


# ── PPT helpers ───────────────────────────────────────────────────────

def _next_slide_no() -> int:
    global _slide_no
    _slide_no += 1
    return _slide_no


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text: str = "Mega-Trend Investment Research 2026") -> None:
    n = _next_slide_no()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(10), Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = SLATE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"  {text}                                                                                    {n}"
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED


def add_accent_line(slide, top=1.05) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(top), Inches(10), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()


def slide_title(slide, title: str, subtitle: str | None = None) -> None:
    set_bg(slide, OFF_WHITE)
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.15))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = NAVY
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.margin_left = Inches(0.45)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = CYAN
    add_accent_line(slide)
    add_footer(slide)


def add_image(slide, path: Path, left=0.45, top=1.35, width=9.1) -> None:
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def add_kpi_row(slide, kpis: list[tuple[str, str, str]], top=1.4) -> None:
    """(value, label, color_hex)"""
    w = 2.15
    for i, (val, label, hex_c) in enumerate(kpis):
        left = 0.45 + i * (w + 0.12)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_LINE
        # top stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(0.08))
        r, g, b = int(hex_c[1:3], 16), int(hex_c[3:5], 16), int(hex_c[5:7], 16)
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = RGBColor(r, g, b)
        stripe.line.fill.background()
        box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.15), Inches(w), Inches(1.1))
        tf = box.text_frame
        p1 = tf.paragraphs[0]
        p1.text = val
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(9)
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.CENTER


def add_bullets(slide, items: list[str], left=0.55, top=1.35, width=4.3, size=13, color=DARK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {annotate(item)}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.line_spacing = 1.2


def add_paragraphs(slide, paragraphs: list[str], left=0.55, top=1.3, width=9.0, size=12, color=DARK, spacing=10):
    """Multi-sentence explanatory body text."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = annotate(para)
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
        p.line_spacing = 1.25


def add_callout(slide, title: str, lines: list[str], left=0.45, top=1.15, width=9.1, accent=CYAN):
    """Key insight / how-to-read box."""
    h = min(2.8, 0.42 + 0.26 * len(lines))
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.08), Inches(h))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(left + 0.22), Inches(top + 0.08), Inches(width - 0.35), Inches(h - 0.12))
    tf = tb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = accent
    for line in lines:
        p = tf.add_paragraph()
        p.text = annotate(line)
        p.font.size = Pt(10)
        p.font.color.rgb = DARK
        p.space_after = Pt(3)


def add_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def chart_slide(prs, title: str, subtitle: str, chart_path: Path, callout_title: str,
                callout_lines: list[str], notes: str, img_top=2.05, img_h=4.5):
    """Chart + top explainer callout."""
    slide = blank_slide(prs)
    slide_title(slide, title, subtitle)
    add_callout(slide, callout_title, callout_lines, top=1.18)
    slide.shapes.add_picture(str(chart_path), Inches(0.45), Inches(img_top), width=Inches(9.1))
    add_notes(slide, notes)
    return slide


def narrative_slide(prs, title: str, subtitle: str, paragraphs: list[str], notes: str = ""):
    slide = blank_slide(prs)
    slide_title(slide, title, subtitle)
    add_paragraphs(slide, paragraphs, size=13, spacing=12)
    if notes:
        add_notes(slide, notes)
    return slide


def section_divider(prs: Presentation, num: str, title: str, subtitle: str, color: RGBColor,
                    overview: list[str] | None = None) -> None:
    slide = blank_slide(prs)
    set_bg(slide, color)
    # decorative circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), Inches(-0.8), Inches(4), Inches(4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(min(color.red + 20, 255) if hasattr(color, 'red') else 0x1E, 0x29, 0x3B)
    circ.fill.transparency = 0.5
    circ.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(8), Inches(3))
    tf = box.text_frame
    p0 = tf.paragraphs[0]
    p0.text = num
    p0.font.size = Pt(56)
    p0.font.bold = True
    p0.font.color.rgb = CYAN
    p1 = tf.add_paragraph()
    p1.text = title
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = MUTED
    p2.space_before = Pt(12)
    if overview:
        p3 = tf.add_paragraph()
        p3.text = ""
        p3.space_before = Pt(18)
        for line in overview:
            p4 = tf.add_paragraph()
            p4.text = f"• {line}"
            p4.font.size = Pt(13)
            p4.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
            p4.space_after = Pt(6)
    add_footer(slide)


def _set_ticker_cell(cell, raw: str) -> None:
    """표 종목 열: 티커(굵게) + 한글명(회색)."""
    if raw in NAME_KO:
        cell.text = NAME_KO[raw]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = DARK
        return
    sym = raw.strip().upper()
    ko = TICKER_KO.get(sym, "")
    if not ko:
        cell.text = raw
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
        return
    cell.text = f"{sym}\n{ko}"
    tf = cell.text_frame
    p1 = tf.paragraphs[0]
    p1.font.bold = True
    p1.font.size = Pt(11)
    p1.font.color.rgb = NAVY
    if len(tf.paragraphs) > 1:
        p2 = tf.paragraphs[1]
        p2.font.size = Pt(10)
        p2.font.color.rgb = MUTED


def add_table(slide, headers, rows, left=0.4, top=1.3, col_widths=None, ticker_cols: set[int] | None = None):
    ticker_cols = ticker_cols or set()
    nr, nc = len(rows) + 1, len(headers)
    row_h = 0.42 if ticker_cols else 0.36
    tbl = slide.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(9.2), Inches(row_h * nr + 0.25)).table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        for p in c.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
    for i, row in enumerate(rows, 1):
        bg = OFF_WHITE if i % 2 else WHITE
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            if j in ticker_cols:
                _set_ticker_cell(c, val)
            elif isinstance(val, str) and any(x in val for x in (" · ", "\n")) and j == 1 and "티커" in headers[1]:
                c.text = val
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.font.color.rgb = DARK
            else:
                c.text = annotate(val) if isinstance(val, str) else str(val)
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.color.rgb = DARK
            c.margin_left = Inches(0.06)
            c.margin_right = Inches(0.04)
            c.margin_top = Inches(0.04)
            c.margin_bottom = Inches(0.04)


def add_native_bar_chart(slide, title: str, categories, series_name, values, left=5.0, top=1.4):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(left), Inches(top), Inches(4.5), Inches(3.8), chart_data
    ).chart
    chart.has_legend = False
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    return chart


# ── Build deck ────────────────────────────────────────────────────────

def build() -> Path:
    global _slide_no
    _slide_no = 0

    # Pre-render charts
    charts = {
        "capex": chart_hyperscaler_capex(),
        "allocation": chart_capex_allocation(),
        "dc_power": chart_dc_power(),
        "market": chart_megatrend_market_2030(),
        "heatmap": chart_demand_heatmap(),
        "glp1": chart_glp1_share(),
        "quantum": chart_quantum_growth(),
        "ionq": chart_ionq_vs_google(),
        "nuclear": chart_nuclear_ppa(),
        "value_chain": chart_value_chain(),
        "horizon": chart_horizon(),
        "stock_count": chart_stock_count(),
    }

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── 1. Cover ──
    slide = blank_slide(prs)
    set_bg(slide, NAVY)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()
    deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(4.5), Inches(5), Inches(5))
    deco.fill.solid()
    deco.fill.fore_color.rgb = SLATE
    deco.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.65), Inches(1.8), Inches(8.5), Inches(4))
    tf = tb.text_frame
    for i, (txt, sz, bold, col) in enumerate([
        ("MEGA-TREND", 18, False, CYAN),
        ("2026–2035\n투자 전략 발표자료", 40, True, WHITE),
        ("8대 메가트렌드 · 수요 시점 · 180+ 종목 분석", 18, False, MUTED),
        ("", 8, False, MUTED),
        ("2026년 7월 5일", 13, False, MUTED),
        ("Morgan Stanley · McKinsey · Gartner · CHIPS Act", 11, False, MUTED),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = col
        if i == 1:
            p.line_spacing = 1.1
    add_footer(slide, "CONFIDENTIAL — For Discussion Purposes")
    add_notes(slide,
        "발표 목적: 2026–2035 주식시장 메가트렌드와 수요 타이밍, 관련 종목을 한 번에 이해할 수 있도록 구성. "
        "투자 권유가 아닌 리서치 자료입니다.")

    # ── 2. TOC ──
    slide = blank_slide(prs)
    slide_title(slide, "목차", "Table of Contents")
    toc = [
        ("01", "Executive Summary", "핵심 메시지 · 시장 규모 · 수요 Heatmap"),
        ("02", "AI 물리 인프라", "2차 파생 — Electrons, Not Tokens"),
        ("03", "원자력 · AI 전력", "PPA Rush · SMR Pipeline"),
        ("04", "반도체 · CHIPS", "리쇼어링 · 48D 마감"),
        ("05", "양자 · PQC", "FTQC Roadmap · 보안 마이그레이션"),
        ("06", "자율 · 로보틱스", "Factory → Humanoid"),
        ("07", "우주 경제", "Launch · Connectivity · Defense"),
        ("08", "바이오 · GLP-1", "Code Meets Cell"),
        ("09", "경제안보", "희토류 · 방산 · Cyber"),
        ("10", "결론 · Appendix", "Horizon · ETF · 면책"),
    ]
    for i, (num, title, sub) in enumerate(toc):
        row, col = divmod(i, 2)
        left = 0.5 + col * 4.8
        top = 1.45 + row * 1.05
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(4.5), Inches(0.88))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_LINE
        box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), Inches(4.2), Inches(0.75))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{num}   {title}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(9)
        p2.font.color.rgb = MUTED

    # ── 2b. How to read ──
    narrative_slide(prs, "이 자료 읽는 법", "발표자 없이도 이해할 수 있도록 구성",
        [
            "각 섹션은 ①왜 중요한지(배경) → ②데이터(차트) → ③어떤 회사가 수혜인지(종목) 순서로 읽으면 됩니다.",
            "파란색 '핵심 해석' 박스는 차트·표를 어떻게 읽어야 하는지, 투자적으로 무엇을 의미하는지 요약합니다.",
            "슬라이드 하단 번호는 순서이며, PowerPoint '발표자 노트'에 보충 설명을 넣어 두었습니다.",
            "★ Rating은 기술·수요·정부 지원·실현 가능성을 종합한 상대 평가이며, 매수 추천이 아닙니다.",
            "종목 표기: NVDA · 엔비디아 형식 — 앞은 티커(미국 주식 코드), 뒤는 한글 회사명입니다.",
        ],
        notes="청중이 자료만 받아도 따라올 수 있게 구조를 안내하는 슬라이드.")

    # ── 2c. Disclaimer ──
    narrative_slide(prs, "면책 및 데이터 출처", "해석 시 유의사항",
        [
            "시가총액·매출·백로그 수치는 2026년 7월 전후 공개 자료 기준이며, 주가 변동에 따라 즉시 달라질 수 있습니다.",
            "수요 시점(예: 2028 FTQC)은 기업 로드맵·애널리스트 추정치이며, 기술·규제 지연 시 2–3년 밀릴 수 있습니다.",
            "출처: Morgan Stanley Big Picture 2026, McKinsey/Gartner Tech Trends, CHIPS Act 공식 문서, Morningstar, IonQ IR, Goldman Sachs AI Capex 리포트 등.",
            "본 자료는 교육·토론 목적이며, 특정 종목 매수·매도를 권유하지 않습니다.",
        ])

    # ── 3. Agenda / Key message ──
    slide = blank_slide(prs)
    slide_title(slide, "오늘의 핵심 논지", "The Great Broadening — 2026")
    add_kpi_row(slide, [
        ("$600B+", "2026 AI Capex", "#06B6D4"),
        ("75–80%", "물리 인프라 비중", "#8B5CF6"),
        ("8", "메가트렌드", "#10B981"),
        ("180+", "분석 종목", "#F59E0B"),
    ])
    add_paragraphs(slide, [
        "2023–25년 주식시장은 'AI = Nvidia·Microsoft 등 소수 메가캡'으로 수익이 집중됐습니다. "
        "2026년부터 Morgan Stanley가 말하는 Great Broadening(수익의 대확산)이 시작됩니다.",
        "이유: AI를 돌리려면 GPU만으로는 부족하고, 전력·변압기·냉각·송전·건물이 먼저 필요합니다. "
        "이 '2차 파생' 산업은 주문이 2–3년치 백로그로 쌓이며, GPU 교체 주기와 무관하게 매출이 이어집니다.",
        "동시에 GLP-1(비만약)·원전 PPA·CHIPS 반도체 리쇼어링은 각각 지금 cash flow가 발생하는 축입니다. "
        "양자·휴머노이드·SMR은 2028년 이후 본격화되는 '다음 파도'입니다.",
    ], top=2.95, size=12, spacing=10)
    add_notes(slide, "청중에게 '왜 메가캡 AI만 보면 안 되는지'를 먼저 각인시키는 슬라이드.")

    # ── SECTION 01 ──
    section_divider(prs, "01", "Executive Summary", "시장 전망 · 수요 타이밍 · 투자 Horizon", NAVY,
        overview=[
            "8대 메가트렌드의 시장 규모와 '언제' 수요가 몰리는지 정리",
            "AI capex가 실제로 어디로 흐르는지 숫자로 확인",
            "투자 기간(0–2년 / 3–5년 / 5–10년)별 우선순위 제시",
        ])

    narrative_slide(prs, "Executive Summary — 한 줄 요약", "2026년 주식시장의 구조 변화",
        [
            "핵심 변화: AI 투자 금액은 커지지만, 그 돈의 대부분은 칩 회사가 아니라 전력·인프라·건설 회사로 갑니다.",
            "전력 병목: 미국 전력망의 70%가 25년 이상 노후입니다. AI 데이터센터는 2030년까지 전력 수요를 2배 가까이 올립니다.",
            "지정학: 중국과의 기술 경쟁으로 반도체·희토류·방산·사이버 보안에 정부 예산이 구조적으로 붙습니다.",
            "인구·건강: GLP-1는 단순 다이어트약이 아니라 당뇨·심혈관·수면 무호흡 등 대사질환 치료의 새 표준으로 자리 잡는 중입니다.",
        ],
        notes="섹션 01 전체를 관통하는 4가지 메시지.")

    chart_slide(prs,
        "하이퍼스케일러 AI Capex 폭발적 성장",
        "Microsoft · Meta · Amazon · Alphabet · Oracle 합산",
        charts["capex"],
        "이 차트가 말하는 것",
        [
            "2024년 $250B → 2027년 $1T: 역사상 가장 가파른 기업 투자 사이클 중 하나",
            "숫자는 '가이던스' 기반 — 이미 발표된 계획이므로 허수가 아님",
            "투자 함의: 이 금액의 75% 이상이 인프라·전력·냉각으로 간다면, NVDA만으로는 부족",
        ],
        "발표 시: '1조 달러가 어디로 가는가'로 다음 슬라이드 연결")

    slide = blank_slide(prs)
    slide_title(slide, "AI Capex는 칩이 아닌 물리 인프라로", "Goldman Sachs: DC 건설·전력·냉각이 capex의 대부분")
    add_callout(slide, "왜 80%가 칩이 아닌가?", [
        "AI 랙 1개(Blackwell 세대)는 약 120kW — 이전 세대(30kW)의 4배 전력",
        "GPU 가격은 떨어질 수 있지만, '전기를 끌어오는 비용'과 '열을 식히는 비용'은 계속 증가",
        "그래서 Vertiv(냉각)·Eaton(변압기)·Caterpillar(건설)가 NVDA와 같이 또는 더 오른 것",
    ], top=1.15)
    add_image(slide, charts["allocation"], left=0.5, top=2.55, width=4.6)
    add_bullets(slide, [
        "Vertiv: 백로그 $15B, 주문이 출하의 2.9배 → 2027년까지 매출 가시성",
        "변압기·스위치기어 납기 최대 5년 → 공급 부족 = 가격 결정력",
        "투자자 관점: 'AI 테마주'가 아닌 '산업 장비주'로 보는 것이 정확",
        "리스크: DC 건설 지연 시 단기 주가 조정, but 백로그는 오히려 늘어남",
    ], left=5.35, top=2.65, width=4.3, size=11)
    add_notes(slide, "2차 파생(Second Derivative) 개념을 반드시 설명할 것.")

    chart_slide(prs,
        "데이터센터 전력 수요",
        "미국 전체 전력 대비 DC 비중",
        charts["dc_power"],
        "해석 가이드",
        [
            "4.4% → 9%: '작아 보이지만' 절대량은 거대 — 신규 발전 설비 없이는 불가능",
            "지역별 전기요금 상승 → 정치·규제 리스크 (유틸리티 투자 시 주의)",
            "원자력·가스·태양+저장이 AI 전력의 3대 축",
        ],
        "전력이 왜 원자력 섹션으로 이어지는지 설명")

    chart_slide(prs,
        "2030 메가트렌드 시장 규모 비교",
        "절대 규모가 큰 순 (로그 스케일 근사)",
        charts["market"],
        "읽는 법",
        [
            "AI Infra(capex)가 압도적 — 단기 자금 유입 최대",
            "GLP-1·Utilities도 수천억 달러 — 이미 상업 매출 발생 중",
            "Quantum·Humanoid는 작지만 성장률(CAGR) 최고 → 장기 옵션",
        ],
        "규모 vs 성장률 트레이드오프 강조")

    chart_slide(prs,
        "메가트렌드별 수요 강도 Heatmap",
        "연도별로 '돈이 몰리는 세기' (0–100 상대 지수)",
        charts["heatmap"],
        "색이 진할수록 = 그 시기에 수요·정책·capex가 집중",
        [
            "AI·바이오·반도체: 2025–26이 이미 정점 구간",
            "양자·로보·우주: 2028 이후 진해짐 → 지금은 '포지셔닝' 단계",
            "100은 '절대 매수 신호'가 아니라 구조적 수요 강도",
        ],
        "Heatmap은 본 보고서의 '수요 시점' 예측의 핵심 도구")

    slide = blank_slide(prs)
    slide_title(slide, "8대 메가트렌드 한눈에", "이름 · 규모 · Peak · 현재 단계")
    add_callout(slide, "표 사용법", [
        "'수요 Peak' = 주가·실적·정책이 동시에 몰리기 쉬운 구간 (정확한 날짜 아님)",
        "'단계' = 백로그(주문잔고) / PPA(전력계약) / Fab(공장) 등 현재 산업 위치",
    ], top=1.12, width=9.0)
    add_table(slide,
        ["#", "트렌드", "2030 규모", "수요 Peak", "단계"],
        [
            ["1", "AI 물리 인프라", "$600B+ capex", "2025–32", "백로그 폭발"],
            ["2", "원자력·전력", "$1.4T 유틸투자", "25–28 / 30–35", "PPA 체결 rush"],
            ["3", "반도체·CHIPS", "$52.7B+", "건설 25–27", "Fab 가동 시작"],
            ["4", "양자+PQC", "$11.5B QC", "PQC 26–31", "조기 상용화"],
            ["5", "로보틱스", "$4–18B", "28–32", "파일럿→양산"],
            ["6", "우주", "$2T(2040)", "26–35", "IPO·발사 cadence"],
            ["7", "GLP-1", "$150–190B", "24–30", "수요>공급"],
            ["8", "경제안보", "$12B+", "25–30", "정책 주도"],
        ], top=1.85, col_widths=[0.35, 1.8, 1.8, 1.6, 1.5])

    chart_slide(prs,
        "투자 Horizon별 섹터 매력도",
        "지금 당장 vs 3년 후 vs 10년 후 — 무엇에 베팅할 것인가",
        charts["horizon"],
        "막대가 높을수록 = 해당 기간에 상대적 투자 매력",
        [
            "0–2년: AI 인프라·GLP-1 — 이미 매출·백로그로 증명",
            "3–5년: 원자력 SMR·양자 — LOI·건설·규제 마일스톤 구간",
            "5–10년: PQC 의무화·바이오 팩토리 — regulation-driven",
        ],
        "개인 투자 기간에 맞춰 섹션을 선택하도록 안내")

    # ── SECTION 02 AI ──
    section_divider(prs, "02", "AI 물리 인프라", "Electrons, Not Tokens · 2025–2032", SLATE,
        overview=[
            "GPU 다음으로 돈이 가는 '전기·냉각·건설' 레이어",
            "3-Layer Stack으로 어떤 회사가 어느 단계인지 구분",
            "백로그·납기가 핵심 KPI",
        ])

    narrative_slide(prs, "AI 물리 인프라란?", "2차 파생(Second Derivative) 투자의 정의",
        [
            "1차 파생: AI 수요가 직접 올리는 회사 — Nvidia(GPU), Microsoft(Azure), OpenAI(API). 토큰·클라우드 사용량에 매출이 연동됩니다.",
            "2차 파생: AI를 '물리적으로 가능하게' 하는 회사 — 전력을 공급하고, 열을 식히고, 건물을 짓습니다. GPU 수요가 늘면 이들 주문도 늘지만, GPU 가격 하락과 무관하게 성장할 수 있습니다.",
            "왜 지금인가: Blackwell 등 차세대 칩은 전력 밀도가 급증해 액체 냉각이 사실상 필수가 됐고, 미국 전력망은 이 속도로 확장되지 않아 병목이 발생합니다.",
            "핵심 지표: book-to-bill(주문/매출 비율), backlog(잔고), lead time(납기) — PER보다 이 숫자가 더 중요한 구간입니다.",
        ])

    chart_slide(prs,
        "AI Value Chain — 자금이 흐르는 곳",
        "하이퍼스케일러 capex → 1차·2차·3차 수혜자",
        charts["value_chain"],
        "다이어그램 읽기",
        [
            "왼쪽: MSFT·META 등이 총 capex 결정",
            "중간 2차(청색): 2026년 주가 상승의 중심 — VRT·ETN·PWR",
            "오른쪽 3차(녹색): 전력 생산 — CEG·VST (원자력 섹션과 연결)",
        ],
        "Value chain을 그린 뒤 각 레이어 종목 슬라이드로")

    slide = blank_slide(prs)
    slide_title(slide, "3-Layer Physical Stack", "칩 → 랙 → 건물·그리드")
    layers = [
        ("Layer 1", "칩 · 인터커넥트",
         f"{ts_block('NVDA','AVGO','MRVL')}\n— GPU·커스텀 ASIC. AI 연산의 핵심.",
         "GPU·ASIC", "#8B5CF6"),
        ("Layer 2", "랙 전력 · 냉각",
         f"{ts_block('VRT','ETN','ABB')}\n— 랙 전력분배·액침냉각. AI DC 매출 비중 급증.",
         "전력·열", "#06B6D4"),
        ("Layer 3", "건물 · 백업 · 그리드",
         f"{ts_block('CAT','PWR','GNRC','STX')}\n— DC 건설·송전·백업발전·데이터 저장.",
         "건설·전력망", "#10B981"),
    ]
    for i, (layer, name, desc, tag, hex_c) in enumerate(layers):
        top = 1.35 + i * 1.95
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(top), Inches(9), Inches(1.75))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_LINE
        r, g, b = int(hex_c[1:3], 16), int(hex_c[3:5], 16), int(hex_c[5:7], 16)
        tag_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(top + 0.2), Inches(1.2), Inches(1.2))
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = RGBColor(r, g, b)
        tag_shape.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.65), Inches(top + 0.45), Inches(1.2), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.text = layer
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tb2 = slide.shapes.add_textbox(Inches(2.0), Inches(top + 0.12), Inches(7.3), Inches(1.45))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p1 = tf2.paragraphs[0]
        p1.text = f"{name}  ({tag})"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        p2 = tf2.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = DARK
    add_notes(slide, "Layer 2가 투자 포커스 — 공급 제약 + 백로그")

    slide = blank_slide(prs)
    slide_title(slide, "핵심 종목 — AI 인프라", "역할 · 재무 포인트 · 리스크")
    add_table(slide,
        ["종목", "하는 일", "왜 지금?", "주의점"],
        [
            ["VRT", "랙 전력·액침냉각", "백로그 $15B, S&P500 편입", "밸류에이션 높음"],
            ["ETN", "변압기·스위치기어", "버티브 대비 저변동, AI 15%+", "순수 AI 플레이 아님"],
            ["PWR", "송전·DC EPC", "전력망 병목 직접 해소", "규제·인허가 지연"],
            ["NVDA", "GPU 풀스택", "AI 1차 수혜, FCF 풍부", "커스텀 ASIC 대체"],
            ["AVGO", "Google/Meta ASIC", "다년 계약·네트워킹", "고객 집중"],
            ["CAT", "DC 건설·발전기", "AI 안 보이는 수혜", "경기 민감성"],
        ], top=1.25, col_widths=[1.1, 1.9, 2.4, 2.0], ticker_cols={0})

    # ── SECTION 03 Nuclear ──
    section_divider(prs, "03", "원자력 · AI 전력", "Nuclear Renaissance · PPA → SMR", RGBColor(0x05, 0x4A, 0x3A),
        overview=[
            "AI DC는 24시간 풀가동 — 태양광만으로는 부족, baseload 전력 필요",
            "지금: 기존 원전 PPA / 나중: SMR 신규 건설",
            "리스크·수익 프로필이 다른 3-Bucket으로 분류",
        ])

    narrative_slide(prs, "왜 AI가 원자력을 찾는가?", "전력 병목의 논리",
        [
            "AI 학습·추론 클러스터는 24/7 가동됩니다. 태양광은 밤에 멈추고, 배터리만으로 GW급을 버티기 어렵습니다.",
            "가스 발전은 가능하지만 탄소·가격 변동 이슈가 있고, 빅테크는 RE100(재생에너지 100%) 목표와 충돌합니다.",
            "원자력은 (1) 연중 무중단 baseload (2) 20년 장기 PPA 가능 (3) GW 단위 확장 — 세 가지를 동시에 만족합니다.",
            "2024–26년 Microsoft–Constellation(TMI 재가동), Meta–6.6GW 패키지 등이 '선례'가 됐고, 이후 발주는 이 선례를 따릅니다.",
        ])

    chart_slide(prs,
        "빅테크 원자력 확보 경쟁",
        "Meta 6.6GW 패키지 구성 (2026.1 발표)",
        charts["nuclear"],
        "차트 설명",
        [
            "녹색 = 이미 돌아가는 원전에서 전력 공급 (CEG, VST) → 매출·현금흐름 즉시",
            "주황 = SMR(소형모듈로) — 2030–35 가동 예정 → 지금은 주가에 '옵션' 반영",
            "투자 시: Bucket 1(유틸)과 Bucket 2(SMR)는 완전히 다른 리스크",
        ],
        "SMR은 pre-revenue — Oklo·NuScale 주의")

    slide = blank_slide(prs)
    slide_title(slide, "투자 3-Bucket Framework", "같은 '원자력 테마' 안에서도 전혀 다른 주식")
    buckets = [
        ("Bucket 1", "Cash-flow 유틸", ts_block("CEG", "VST", "TLN", "D"),
         "이미 돌아가는 원전 + 20년 PPA. AI DC 전력을 '팔고' 있는 상태. 배당·FCF 가능.",
         "중간", "#10B981"),
        ("Bucket 2", "SMR Pure-play", ts_block("OKLO", "SMR", "NNE"),
         "아직 매출 거의 없음. NRC 인허가·건설·고객 LOI가 핵심. 5–10년 binary 옵션.",
         "매우 높음", "#F59E0B"),
        ("Bucket 3", "연료·인프라", ts_block("CCJ", "BWXT", "LEU", "GEV"),
         "우라늄·핵연료·터빈. 원전 대수가 늘면 같이 성장. 방산 수요도 겹침.",
         "중간", "#06B6D4"),
    ]
    for i, (b, title, tickers, desc, risk, hex_c) in enumerate(buckets):
        top = 1.3 + i * 1.9
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(top), Inches(9), Inches(1.65))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_LINE
        r, g, b = int(hex_c[1:3], 16), int(hex_c[3:5], 16), int(hex_c[5:7], 16)
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(top + 0.12), Inches(8.5), Inches(1.45))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{b}  |  {title}  |  리스크: {risk}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(r, g, b)
        p2 = tf.add_paragraph()
        p2.text = tickers
        p2.font.size = Pt(10)
        p2.font.bold = False
        p2.font.color.rgb = NAVY
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = DARK

    # ── SECTION 04 Semi ──
    section_divider(prs, "04", "반도체 · CHIPS", "Reshoring · 48D Cliff Dec 2026", RGBColor(0x3B, 0x1D, 0x6E),
        overview=[
            "미국 advanced logic 점유율 12% → 22% 회복 중",
            "2026.12.31 48D 세액공제 착공 마감 = 긴급 건설 rush",
            "Fab vs 장비사 — 수혜 구조가 다름",
        ])

    narrative_slide(prs, "CHIPS Act를 이해해야 하는 이유", "지정학 + AI = 반도체 리쇼어링",
        [
            "첨단 칩의 90% 이상이 대만(TSMC)에 집중돼 있다는 것은 미국에게 국가 안보 리스크입니다. 2022년 CHIPS Act는 이를 해소하기 위한 $52.7B+ 패키지입니다.",
            "직접 지원: Intel $8.5B, TSMC $6.6B, Samsung, Micron 등 — fab 건설·현대화.",
            "48D 세액공제: 건설 투자의 35%를 세금으로 돌려받음. 단, 2026년 12월 31일 이전 착공해야 함 → 지금이 건설 피크.",
            "투자 포인트: (A) Fab 운영사 Intel·TSM — turnaround/성장 (B) 장비사 AMAT·LRCX·ASML — fab가 늘면 주문이 늘음, CHIPS 의존도 낮음.",
        ])

    slide = blank_slide(prs)
    slide_title(slide, "CHIPS Act 현황 — 숫자로 보기", "2026년 7월 기준")
    add_kpi_row(slide, [
        ("$38.7B", "/39B 배정 완료", "#8B5CF6"),
        ("35%", "48D Tax Credit", "#06B6D4"),
        ("12/31/26", "착공 Dead-line", "#F43F5E"),
        ("22%", "US Adv Logic", "#10B981"),
    ], top=1.25)
    add_paragraphs(slide, [
        "Intel Ohio는 2030으로 지연됐지만, '착공'만 2026 안에 하면 세액공제는 유지 — 실제 wafer 출하는 몇 년 뒤일 수 있음.",
        "Intel 18A(Arizona)는 이미 HVM(대량생산) — 미국 내 2nm급 파운드리의 상징적 마일스톤.",
        "Secure Enclave: Intel에 DoD $3B — 군·정보기관용 보안 칩 공급망.",
    ], top=2.85, size=12)

    slide = blank_slide(prs)
    slide_title(slide, "CHIPS 주요 수혜 기업", "누가 얼마를 받았고, 무엇을 짓는가")
    add_table(slide,
        ["종목", "지원 규모", "프로젝트", "투자자 관점"],
        [
            ["INTC", "$8.5B+대출", "AZ·OH·OR 18A", "턴어라운드, 고위험"],
            ["TSM", "$6.6B", "Arizona Fab21", "품질·실행력 최상"],
            ["Samsung", "$6.4B", "Taylor 2nm", "메모리+로직"],
            ["MU", "$6.16B", "NY·Idaho DRAM", "AI HBM 수혜"],
            ["AMAT", "—(간접)", "장비 1위", "Capex cycle 레버리지"],
            ["ASML", "—(간접)", "EUV 독점", "필수 인프라"],
        ], top=1.25, col_widths=[1.1, 1.2, 2.2, 2.4], ticker_cols={0})

    # ── SECTION 05 Quantum ──
    section_divider(prs, "05", "양자 · PQC", "FTQC 2028–30 · Security Migration", RGBColor(0x78, 0x35, 0x0F),
        overview=[
            "양자컴퓨팅: 아직 초기 상용화, but 로드맵이 2028–30으로 수렴",
            "PQC(포스트퀀텀 암호): 지금 당장 마이그레이션 시작 — regulation driven",
            "하드웨어 vs 보안 — 수익 시점이 다름",
        ])

    narrative_slide(prs, "양자컴퓨팅 + PQC — 왜 둘 다 중요한가?", "공격 기술과 방어 기술이 동시에 움직임",
        [
            "양자컴퓨터가 충분히 강력해지면, 오늘 RSA로 암호화된 데이터를 나중에 해독할 수 있습니다(지금 훔쳐 두고 나중에 푼다 = Harvest Now, Decrypt Later).",
            "그래서 NIST가 2024년 PQC 표준(FIPS 203/204/205)을 확정했고, 미국 정부는 2030–31까지 연방기관·계약업체의 마이그레이션을 요구합니다.",
            "양자 '하드웨어' 투자(IonQ 등)는 2028 FTQC(장애허용 양자컴퓨터) 전후가 관건이고, PQC 투자(PANW, CRWD 등)는 지금부터 5년간 매출이 발생합니다.",
            "결론: 단기 현금흐름은 PQC·클라우드 접근(QaaS)이, 장기 옵션은 하드웨어 pure-play가 큽니다.",
        ])

    chart_slide(prs,
        "양자컴퓨팅 시장 성장",
        "하드웨어 + 클라우드 + 서비스 합산",
        charts["quantum"],
        "시장 단계",
        [
            "2025–27: 파일럿·클라우드 접근·정부 R&D — 매출은 작지만 CAGR 40%+",
            "2028–30: FTQC 목표 연도 — 로드맵 달성 여부가 주가 분기점",
            "PQC 시장은 양자 시장보다 절대 규모가 클 수 있음(엔터프라이즈 전체 마이그레이션)",
        ],
        "IonQ Q1'26 매출 $64.7M — pure-play 중 유일한 상업 스케일")

    slide = blank_slide(prs)
    slide_title(slide, "IonQ vs Google — 누가 앞서는가?", "기술 경로가 다르면 '우위'도 다름")
    add_callout(slide, "비교 전제", [
        "IonQ = 상장 pure-play, 매출·RPO로 검증 가능",
        "Google = Alphabet 내부, 연구·QEC에서 선행, 상용 매출은 작음",
        "점수는 상대 비교(0–100)이며 절대 기술 등급이 아님",
    ], top=1.12)
    add_image(slide, charts["ionq"], top=2.35, width=9.0)
    add_notes(slide,
        "IonQ: 트랩드 이온, 게이트 충실도 99.99%, all-to-all 연결, Q1'26 RPO $470M. "
        "Google Willow: below-threshold QEC 2024 최초 실증, 큐비트 수·확장성 우위. "
        "투자: IonQ=상용 베팅, GOOGL=간접 R&D 옵션.")

    slide = blank_slide(prs)
    slide_title(slide, "양자 + PQC 종목 맵", "구분별로 수익 시점이 다름")
    add_table(slide,
        ["구분", "대표 종목", "무엇을 파는가", "언제 돈이 되는가"],
        [
            ["하드웨어", ts_cell("IONQ RGTI QBTS"), "양자 프로세서·클라우드", "2028 FTQC 전후"],
            ["빅테크", ts_cell("IBM GOOGL MSFT"), "QaaS·연구·생태계", "지금~지속"],
            ["PQC 소프트", ts_cell("PANW CRWD NET ZS"), "TLS·방화벽·제로트러스트", "2026–31 rush"],
            ["PQC 하드웨어", ts_cell("LAES WKEY NXPI"), "보안 칩·HSM", "2027 인증 의무"],
        ], top=1.25, col_widths=[1.0, 2.6, 2.3, 2.0])

    # ── SECTION 06 Robotics ──
    section_divider(prs, "06", "자율 · 로보틱스", "Factory Floor → Humanoid", RGBColor(0x7F, 0x1D, 0x1D),
        overview=[
            "AI가 '분석'에서 '행동'으로 — Physical AI",
            "지금: 창고·공장 / 다음: 휴머노이드",
            "비상장 Figure AI가 valuation 선도",
        ])

    narrative_slide(prs, "Autonomy Stack이란?", "Jensen Huang: '움직이는 모든 것이 자율화된다'",
        [
            "1단계(지금): 창고·공장 — Amazon 75만 대+, Symbotic(Walmart), Fanuc·ABB 산업로봇에 AI 비전·경로계획 탑재.",
            "2단계(2026–28): 물류·검사·위험 작업 — 드론(AVAV), 수술로봇(ISRG)은 이미 수익화.",
            "3단계(2028–32): 휴머노이드 — Tesla Optimus, Figure AI, Boston Dynamics Atlas. 아직 '매출'보다 '파일럿·생산속도'가 뉴스.",
            "투자 현실: pure-play 휴머노이드 상장사는 거의 없음(TSLA, UBTECH). 부품·SI(NVDA, TER, CGNX)가 안전한 간접 노출.",
        ])

    slide = blank_slide(prs)
    slide_title(slide, "로보틱스 — 지표와 종목", "2026년 확인해야 할 마일스톤")
    add_kpi_row(slide, [
        ("750K+", "Amazon Robots", "#F43F5E"),
        ("1대/시간", "Figure BotQ", "#06B6D4"),
        ("37.9%", "CAGR 26-30", "#10B981"),
        ("$39B", "Figure AI Val", "#8B5CF6"),
    ], top=1.25)
    add_paragraphs(slide, [
        "Tesla Optimus: 공장 내 테스트 중이나 Musk 본인도 '아직 material usage 아님' 인정 — 2027 consumer 판매가 re-rating 트리거.",
        "Figure AI: 비상장, BotQ 공장 1대/시간 생산 — valuation $39B는 '기대'가 반영된 수치.",
        "확인할 것: 외부 유료 고객 계약 공개, 단가·내구시간·작업 성공률 — 데모 영상만으로는 부족.",
    ], top=2.85, size=12)

    # ── SECTION 07 Space ──
    section_divider(prs, "07", "우주 경제", "SpaceX IPO · Launch Economy", RGBColor(0x1E, 0x3A, 0x8A),
        overview=[
            "발사 비용 하락 → 위성·국방·통신 수요 폭발",
            "SpaceX IPO = 섹터 벤치마크 탄생",
            "Launch / Connect / Observe / Defense 4축",
        ])

    narrative_slide(prs, "우주 경제가 커지는 이유", "더 싸게, 더 자주, 더 많이 올린다",
        [
            "SpaceX가 Falcon 9 재사용으로 kg당 발사 비용을 10배 이상 낮췄습니다. '우주는 정부 전용'에서 '민간 인프라'로 바뀌는 전환점입니다.",
            "2026년 SpaceX 상장(SPCX)은 우주 경제에 '가격 표'가 생긴 것 — 다른 우주주는 SpaceX 대비 할인·프리미엄으로 거래됩니다.",
            "수요 4대 축: (1) Starlink류 위성통신 (2) Earth Observation AI (3) 국방·정찰 (4) 달·궤도 서비스.",
            "Rocket Lab(RKLB): 미국 2위 발사, Neutron으로 중형 시장 진입 시도 — SpaceX 다음으로 순수 플레이에 가까움.",
        ])

    slide = blank_slide(prs)
    slide_title(slide, "우주 경제 — 4축 투자 맵", "각 축의 역할과 대표 종목")
    cols = [
        ("Launch · 발사", ts_block("SPCX", "RKLB", "FLY"),
         "로켓·엔진. 발사 횟수·성공률·단가가 KPI.", "#3B82F6"),
        ("Connectivity · 통신", ts_block("ASTS", "IRDM", "GSAT"),
         "우주에서 스마트폰·IoT 직연결. 스펙트럼·위성 수가 핵심.", "#8B5CF6"),
        ("Earth Obs · 관측", ts_block("PL", "BKSY"),
         "위성 이미지 + AI 분석. 농업·보험·군사.", "#06B6D4"),
        ("Defense · 방산", ts_block("LMT", "NOC", "RTX", "LHX"),
         "위성·미사일·전자전. 정부 예산 연동.", "#10B981"),
    ]
    for i, (title, tickers, desc, hex_c) in enumerate(cols):
        row, col = divmod(i, 2)
        left = 0.5 + col * 4.7
        top = 1.35 + row * 2.75
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(4.4), Inches(2.45))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_LINE
        r, g, b = int(hex_c[1:3], 16), int(hex_c[3:5], 16), int(hex_c[5:7], 16)
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(4.4), Inches(0.5))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = RGBColor(r, g, b)
        hdr.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(left), Inches(top + 0.06), Inches(4.4), Inches(0.45))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tb2 = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.58), Inches(4.1), Inches(1.7))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = tickers
        p2.font.size = Pt(9)
        p2.font.bold = False
        p2.font.color.rgb = NAVY
        p3 = tf2.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(9)
        p3.font.color.rgb = DARK

    # ── SECTION 08 Bio ──
    section_divider(prs, "08", "바이오 · GLP-1", "Code Meets Cell · $150–190B", RGBColor(0x83, 0x15, 0x5C),
        overview=[
            "비만약 → 대사질환 표준요법",
            "LLY·NVO duopoly, 2030까지 93%",
            "다음: oral·triple agonist·근손실 개선",
        ])

    narrative_slide(prs, "GLP-1 — 왜 메가트렌드인가?", "Morgan Stanley 'Code Meets Cell'",
        [
            "GLP-1(세마글루타이드·티르제파타이드)는 체중 감량뿐 아니라 당뇨, 심혈관 위험 감소, 수면무호흡 등으로 적응증이 확장 중입니다.",
            "2030년 시장 $150–200B, 환자 5,900만 명(TD Cowen) — 역대 단일 약물 클래스 중 최대 규모 후보.",
            "LLY 62% + NVO 31% = 사실상 duopoly. 신규 진입자(AMGN, VKTX)는 '더 나은 효능·oral·근육 보존'으로 틈새를 노림.",
            "리스크: Medicare 가격 협상, 공급 과잉(2028+), oral·제네릭으로 가격 하락 — but volume이 compensating.",
        ])

    slide = blank_slide(prs)
    slide_title(slide, "GLP-1 시장 구조", "누가 얼마나 가져가는가")
    add_callout(slide, "차트 읽기", [
        "2030년 duopoly 93% — 신규 진입이 어렵지만, oral·triple agonist가 game changer 될 수 있음",
        "투자: LLY/NVO=core · AMGN/VKTX=optionality · TMO/DHR=피킹스앤샤벨",
    ], top=1.12, width=4.8, left=0.45)
    add_image(slide, charts["glp1"], left=0.5, top=2.2, width=4.5)
    add_bullets(slide, [
        "LLY: Mounjaro/Zepbound + orforglipron(oral) + retatrutide(triple)",
        "NVO: Ozempic/Wegovy + CagriSema(amylin combo)",
        "다음 승자 조건: (1) oral 순응도 (2) 근손실 적음 (3) 심장·간 데이터",
        "공급: 2026부터 점진적 완화, but 여전히 병목 구간",
    ], left=5.2, top=2.3, width=4.5, size=11)

    # ── SECTION 09 Security ──
    section_divider(prs, "09", "경제안보", "Critical Minerals · Defense · Cyber", RGBColor(0x0F, 0x4C, 0x4A),
        overview=[
            "중국 의존 탈피 = 미국 정책 최우선",
            "희토류·구리·방산·사이버 4축",
            "MP Materials = 희토류의 대표 수혜",
        ])

    narrative_slide(prs, "경제안보 투자란?", "Multipolar World — Morgan Stanley 테마",
        [
            "미·중 기술 경쟁으로 '공급망을 자국에 둔다'는 정책이 CHIPS, IRA, DPA, Project Vault($12B)로 구체화됩니다.",
            "희토류: NdFeB 자석은 EV·풍력·미사일·AI 서버에 필수. 중국이 정제의 대부분을 장악 → MP Materials가 미국 유일 통합 체인.",
            "구리: 전력망·DC·EV 모두 구리 소비 증가. FCX 등 미국/칠레 광산.",
            "방산: 우크라·중동·대만 긴장으로 예산 구조적 증가. LMT·NOC·RTX.",
            "사이버: PQC 마이그레이션 + AI 공격 증가 → CRWD·PANW.",
        ])

    slide = blank_slide(prs)
    slide_title(slide, "경제안보 — 정책과 종목", "정부 돈이 따라오는 곳")
    add_kpi_row(slide, [
        ("$12B", "Project Vault", "#14B8A6"),
        ("15%", "DoD→MP 지분", "#06B6D4"),
        ("$110/kg", "NdPr 가격 바닥", "#F59E0B"),
        ("43+", "Critical Min", "#8B5CF6"),
    ], top=1.25)
    add_paragraphs(slide, [
        "MP Materials: DoD 15% 지분 + NdPr $110/kg 10년 floor → 중국 수출 중단 후 미국 내 magnet 생산(2025.12~).",
        "PLTR: 정부·군 AI 소프트웨어 — '실적 모멘텀 + 정책' 겹침.",
        "주의: 정책 테마주는 선거·행정 변경 시 변동성 극대 — 계약·지급 실적을 분기별 확인.",
    ], top=2.85, size=12)

    # ── SECTION 10 Conclusion ──
    section_divider(prs, "10", "결론 · Appendix", "Action Framework", NAVY,
        overview=[
            "기간별 액션 플랜",
            "ETF·리스크 체크리스트",
            "면책",
        ])

    chart_slide(prs,
        "트렌드별 커버리지",
        "본 보고서에서 분석한 종목 수 (각 20+)",
        charts["stock_count"],
        "의미",
        [
            "종목 수가 많다 = 테마 내 분산·세부 niches 존재",
            "ETF로 broad exposure, 개별주로 conviction bet",
        ],
        "종목 리스트 전체는 채팅 보고서 원문 참조")

    slide = blank_slide(prs)
    slide_title(slide, "투자 실행 프레임워크", "내 기간에 맞는 섹션을 고른다")
    add_table(slide,
        ["투자 기간", "무엇을 볼 것인가", "왜?", "ETF 예시"],
        [
            ["0–2년", ts_cell("VRT ETN CEG LLY AMAT"), "백로그·PPA·매출 증명", "SOXX · NUKZ"],
            ["3–5년", ts_cell("OKLO IONQ RKLB TSLA"), "마일스톤·로드맵", "ARKX · BOTZ"],
            ["5–10년", "PQC · 유전자편집 · 달 탐사", "규제 deadline", "HACK · XBI"],
        ], top=1.3, col_widths=[0.9, 3.2, 2.3, 1.3])
    add_callout(slide, "리스크 체크리스트", [
        "□ 밸류에이션: AI 인프라주는 이미 많이 올랐는지 PER·backlog 대비 확인",
        "□ 실행 리스크: SMR·휴머노이드·양자는 '연기'가 일상적",
        "□ 정치 리스크: 유틸 요금·약가 협상·CHIPS 예산 삭감",
        "□ 분산: 한 테마에 올인보다 Horizon별 버킷 구성",
    ], top=3.5)

    narrative_slide(prs, "최종 정리", "한 페이지로 돌아보기",
        [
            "2026–30년 주식시장의 중심축은 'AI가 만든 전기·실물·안보 수요'입니다. GPU 한 종목이 아니라, 전기를 공급하고(CEG), 식히고(VRT), 반도체를 만들고(AMAT), 약을 팔고(LLY) 하는 회사들이 같이 움직입니다.",
            "지금 당장 수요가 몰린 곳: AI 2차 인프라, GLP-1, 운영 원전 PPA, CHIPS fab·장비.",
            "3–5년 후 본격화: SMR 가동, FTQC, 휴머노이드 B2B, 우주 Neutron.",
            "5–10년 regulation play: PQC 전면 의무화, 바이오 팩토리, 달 경제.",
            "Great Broadening — 좁은 승자에서 넓은 기회로. 본 자료는 투자 권유가 아닙니다.",
        ])

    # ── Thank you ──
    slide = blank_slide(prs)
    set_bg(slide, NAVY)
    accent2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.4), Inches(10), Inches(0.08))
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = CYAN
    accent2.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(8.6), Inches(2.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "Q & A"
    p2.font.size = Pt(24)
    p2.font.color.rgb = CYAN
    p2.space_before = Pt(16)
    p3 = tf.add_paragraph()
    p3.text = "Great Broadening — 좁은 승자에서 넓은 기회로"
    p3.font.size = Pt(14)
    p3.font.color.rgb = MUTED
    p3.space_before = Pt(24)
    add_footer(slide)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    target = OUT
    try:
        prs.save(str(target))
    except PermissionError:
        target = OUT_FALLBACK
        prs.save(str(target))
    return target


if __name__ == "__main__":
    path = build()
    print(f"Saved: {path} ({_slide_no} slides)")
